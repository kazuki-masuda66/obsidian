import sys
import os
from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        result = []
        
        result.append(f"=== PPTX File: {os.path.basename(pptx_path)} ===")
        result.append(f"Total Slides: {len(prs.slides)}")
        result.append("")
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            result.append(f"--- Slide {slide_idx} ---")
            
            # スライドのタイトルを取得
            if slide.shapes.title:
                result.append(f"Title: {slide.shapes.title.text}")
            
            # 全てのシェイプからテキストを抽出
            for shape_idx, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        result.append(f"Shape {shape_idx}: {text}")
                
                # テーブルの場合
                if shape.has_table:
                    table = shape.table
                    result.append(f"Table {shape_idx}:")
                    for row_idx, row in enumerate(table.rows, start=1):
                        row_data = []
                        for cell in row.cells:
                            if cell.text:
                                row_data.append(cell.text.strip())
                        if row_data:
                            result.append(f"  Row {row_idx}: {' | '.join(row_data)}")
            
            result.append("")
        
        return "\n".join(result)
    except Exception as e:
        import traceback
        return f"Error extracting text from PPTX: {e}\n{traceback.format_exc()}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <file_path>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = extract_text_from_pptx(file_path)
    
    # Ensure output is encoded in UTF-8 to handle special characters
    sys.stdout.buffer.write(result.encode('utf-8', errors='replace'))


